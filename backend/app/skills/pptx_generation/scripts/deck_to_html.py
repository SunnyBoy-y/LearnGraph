#!/usr/bin/env python3
"""Convert a PPTX's text content into a printable, self-contained HTML preview.

Offline alternative to rendering PPTX to PDF/PNG (no LibreOffice in the image):
emit styled HTML that document-conversion/html_to_pdf or html_to_png can render.
"""
from __future__ import annotations

import argparse
import hashlib
import html as _html
import json
import sys
from pathlib import Path


def _safe(value: str) -> Path:
    path = Path(value)
    if path.is_absolute() or ".." in path.parts:
        raise RuntimeError("path must stay inside the sandbox workspace")
    return path


_CSS = (
    "body{font-family:'Noto Sans CJK SC',sans-serif;margin:2.5em;line-height:1.6;}"
    ".slide{border:1px solid #ccc;border-radius:8px;padding:1.4em 1.8em;margin:0 0 1.6em;"
    "page-break-inside:avoid;}"
    ".slide h2{color:#444;border-bottom:2px solid #4472C4;padding-bottom:.3em;font-size:1.5em;}"
    ".slide ul{margin:.6em 0 0;padding-left:1.3em;}.slide li{margin:.25em 0;}"
    ".notes{color:#777;font-size:.85em;margin-top:.8em;}"
)


def to_html(src: Path, deck_title: str = "") -> str:
    from pptx import Presentation

    prs = Presentation(str(src))
    title = deck_title or (prs.core_properties.title or "PPTX 预览")
    body_parts: list[str] = []
    for slide in prs.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        texts.append(text)
        notes = ""
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        heading = _html.escape(texts[0]) if texts else "（无标题）"
        items = "\n".join(
            f"<li>{_html.escape(item)}</li>" for item in texts[1:]
        )
        notes_html = f'<div class="notes">备注：{_html.escape(notes)}</div>' if notes else ""
        body_parts.append(
            f'<div class="slide"><h2>{heading}</h2>'
            f'<ul>{items}</ul>{notes_html}</div>'
        )
    return (
        "<!doctype html><html><head><meta charset='utf-8'>"
        f"<title>{_html.escape(title)}</title><style>{_CSS}</style></head>"
        f"<body><h1>{_html.escape(title)}</h1>{''.join(body_parts)}</body></html>"
    )


def main() -> int:
    parser = argparse.ArgumentParser(description="PPTX text -> printable HTML preview (offline).")
    parser.add_argument("--input", required=True, help="workspace-relative .pptx path")
    parser.add_argument("--output", required=True, help="workspace-relative .html path")
    parser.add_argument("--title", default="", help="optional deck title override")
    parser.add_argument("--overwrite", action="store_true")
    args = parser.parse_args()
    src = _safe(args.input)
    dst = _safe(args.output)
    if not src.is_file():
        raise RuntimeError(f"input file not found: {src}")
    if dst.exists() and not args.overwrite:
        raise RuntimeError("output already exists; pass --overwrite to replace")
    html = to_html(src, args.title)
    if not html.strip():
        raise RuntimeError("pptx produced empty HTML")
    dst.parent.mkdir(parents=True, exist_ok=True)
    dst.write_text(html, encoding="utf-8")
    print(
        json.dumps(
            {
                "status": "ok",
                "input": str(src),
                "output": str(dst),
                "chars": len(html),
                "sha256": hashlib.sha256(html.encode("utf-8")).hexdigest(),
            },
            ensure_ascii=False,
        )
    )
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except RuntimeError as exc:
        print(json.dumps({"status": "error", "error": str(exc)}, ensure_ascii=False), file=sys.stderr)
        raise SystemExit(1) from exc
