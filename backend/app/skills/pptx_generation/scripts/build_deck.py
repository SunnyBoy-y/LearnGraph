#!/usr/bin/env python3
"""Build a PPTX deck from a structured outline JSON (offline, python-pptx)."""
from __future__ import annotations

import argparse
import hashlib
import json
import sys
from pathlib import Path


def _safe(value: str) -> Path:
    path = Path(value)
    if path.is_absolute() or ".." in path.parts:
        raise RuntimeError("path must stay inside the sandbox workspace")
    return path


def _hex_to_rgb(value: str):
    from pptx.dml.color import RGBColor

    value = (value or "").lstrip("#")
    if not (len(value) == 6 and all(c in "0123456789abcdefABCDEF" for c in value)):
        raise RuntimeError(f"invalid accent color: {value!r}")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def build(outline: dict, dst: Path) -> int:
    from pptx import Presentation
    from pptx.util import Pt

    slides_spec = outline.get("slides")
    if not isinstance(slides_spec, list) or not slides_spec:
        raise RuntimeError("outline.slides must be a non-empty list")
    if len(slides_spec) > 100:
        raise RuntimeError("outline.slides exceeds 100 pages")

    theme = outline.get("theme") or {}
    accent = _hex_to_rgb(theme.get("accent", "4472C4"))
    title_size = int(theme.get("title_size", 32))
    point_size = int(theme.get("point_size", 18))

    prs = Presentation()
    # 16:9 widescreen.
    prs.slide_width = 13_333_200
    prs.slide_height = 7_500_000
    blank = prs.slide_layouts[6]

    for slide_spec in slides_spec:
        if not isinstance(slide_spec, dict) or not str(slide_spec.get("title") or "").strip():
            raise RuntimeError("each slide needs a non-empty title")
        slide = prs.slides.add_slide(blank)
        # Accent rule under the title.
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(1.1))
        tf = title_box.text_frame
        tf.text = str(slide_spec["title"])
        tf.paragraphs[0].font.size = Pt(title_size)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = accent

        points = slide_spec.get("points") or []
        body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(12.0), Inches(5.0))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        for i, point in enumerate(points):
            p = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
            p.text = str(point)
            p.font.size = Pt(point_size)
            p.space_after = Pt(6)
        notes = str(slide_spec.get("notes") or "").strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    dst.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dst))
    return len(slides_spec)


def main() -> int:
    parser = argparse.ArgumentParser(description="Build a PPTX deck from outline JSON (offline).")
    parser.add_argument("--input", required=True, help="workspace-relative outline.json path")
    parser.add_argument("--output", required=True, help="workspace-relative .pptx path")
    parser.add_argument("--overwrite", action="store_true")
    args = parser.parse_args()
    src = _safe(args.input)
    dst = _safe(args.output)
    if not src.is_file():
        raise RuntimeError(f"input file not found: {src}")
    if dst.exists() and not args.overwrite:
        raise RuntimeError("output already exists; pass --overwrite to replace")
    try:
        outline = json.loads(src.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise RuntimeError(f"invalid outline JSON: {exc}") from exc
    if not isinstance(outline, dict):
        raise RuntimeError("outline must be a JSON object")
    slide_count = build(outline, dst)
    data = dst.read_bytes()
    print(
        json.dumps(
            {
                "status": "ok",
                "input": str(src),
                "output": str(dst),
                "slides": slide_count,
                "bytes": len(data),
                "sha256": hashlib.sha256(data).hexdigest(),
            },
            ensure_ascii=False,
        )
    )
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except RuntimeError as exc:
        print(json.dumps({"status": "error", "error": str(exc)}, ensure_ascii=False), file=sys.stderr)
        raise SystemExit(1) from exc
