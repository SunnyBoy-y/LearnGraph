#!/usr/bin/env python3
"""Inspect a PPTX: slide count, per-slide text, and shape statistics (offline)."""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path


def _safe(value: str) -> Path:
    path = Path(value)
    if path.is_absolute() or ".." in path.parts:
        raise RuntimeError("path must stay inside the sandbox workspace")
    return path


def inspect(src: Path) -> dict:
    from pptx import Presentation

    prs = Presentation(str(src))
    slides: list[dict] = []
    for index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        shapes = 0
        for shape in slide.shapes:
            shapes += 1
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        texts.append(text)
        notes = ""
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        slides.append(
            {
                "index": index,
                "title": texts[0] if texts else "",
                "texts": texts,
                "shape_count": shapes,
                "notes": notes,
            }
        )
    return {"slide_count": len(slides), "slides": slides}


def main() -> int:
    parser = argparse.ArgumentParser(description="Inspect a PPTX (offline).")
    parser.add_argument("--input", required=True, help="workspace-relative .pptx path")
    parser.add_argument("--output", default="", help="optional workspace-relative .json path for the report")
    parser.add_argument("--overwrite", action="store_true")
    args = parser.parse_args()
    src = _safe(args.input)
    if not src.is_file():
        raise RuntimeError(f"input file not found: {src}")
    report = inspect(src)
    report.update({"status": "ok", "input": str(src)})
    if args.output:
        dst = _safe(args.output)
        if dst.exists() and not args.overwrite:
            raise RuntimeError("output already exists; pass --overwrite to replace")
        dst.parent.mkdir(parents=True, exist_ok=True)
        dst.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(report, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except RuntimeError as exc:
        print(json.dumps({"status": "error", "error": str(exc)}, ensure_ascii=False), file=sys.stderr)
        raise SystemExit(1) from exc
